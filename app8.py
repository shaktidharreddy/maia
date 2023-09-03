#enhancing -current working version, modifying the pubmed chatbot with pandas df agent
import os
import re
import urllib
import urllib.request
import base64
import tempfile
import json
import nltk
nltk.download('all')
import streamlit as st
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from docx import Document
from docx.shared import Pt
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from datetime import datetime  # Import the 'datetime' class from the 'datetime' module
import requests
from Bio import Entrez
from Bio import Medline
import time
from io import BytesIO
from PIL import Image
from langchain.document_loaders import PyPDFLoader
from langchain.memory import ConversationBufferMemory
from langchain.memory.chat_message_histories import StreamlitChatMessageHistory
from langchain.embeddings import HuggingFaceEmbeddings
from langchain.callbacks.base import BaseCallbackHandler
from langchain.chains import ConversationalRetrievalChain
from langchain.vectorstores import DocArrayInMemorySearch
from langchain.text_splitter import RecursiveCharacterTextSplitter
from shakti_stream_index import llama_vector_index
from streamlit_pills import pills
import streamlit_authenticator as stauth
from streamlit_option_menu import option_menu
from PyPDF2 import PdfReader
from langchain.text_splitter import CharacterTextSplitter
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain.vectorstores import FAISS
from langchain.chains.question_answering import load_qa_chain
from langchain.llms import OpenAI
from langchain.callbacks import get_openai_callback
from langchain.agents import create_json_agent, AgentExecutor
from langchain.agents import create_pandas_dataframe_agent
from langchain.agents.agent_types import AgentType
from streamlit_chat import message
from langchain.agents.agent_toolkits import JsonToolkit
from langchain.chains import LLMChain
from langchain.chat_models import ChatOpenAI
from langchain.requests import TextRequestsWrapper
from langchain.tools.json.tool import JsonSpec
from streamlit_image_select import image_select
from streamlit_chat import message
from llama_index import GPTVectorStoreIndex, SimpleDirectoryReader
from llama_index import download_loader, StorageContext, load_index_from_storage
from vaderSentiment.vaderSentiment import SentimentIntensityAnalyzer
from langchain.callbacks.streamlit import StreamlitCallbackHandler
from wordcloud import WordCloud
import pandas as pd
import matplotlib.pyplot as plt
import numpy as np
#load inthe NTLK stopwords to remove articles, preposition and other words that are not actionable
from nltk.corpus import stopwords
# This allows to create individual objects from a bog of words
from nltk.tokenize import word_tokenize
# Lemmatizer helps to reduce words to the base form
from nltk.stem import WordNetLemmatizer
# Ngrams allows to group words in common pairs or trigrams..etc
from nltk import ngrams
# We can use counter to count the objects
from collections import Counter
# This is our word freq distribution library
from nltk import FreqDist
import seaborn as sns
from time import sleep
from stqdm import stqdm
import itertools
import pickle
import glob

saved_path = os.path.abspath("./") + "/bot_data"
rootdir = os.path.abspath("./")
datadir = os.path.abspath("./") + "/data"
promptdir = os.path.abspath("./") + "/prompts"
Entrez.email = "shakti20889@gmail.com"

# def progress_bar_method(secs):
#     # Code for your second asynchronous method goes here
#     for i in stqdm(range(secs), backend=True, frontend=True):
#         sleep(0.5)

def generate_response1(input_text, df):
    agent = create_pandas_dataframe_agent(ChatOpenAI(temperature =0, model_name="gpt-4", streaming = True), df, verbose=False)
    query_response = agent.run(input_text)
    return query_response


def search_pubmed(article_title, retmax=5):
    # Perform the PubMed search using the article title
    handle = Entrez.esearch(db="pubmed", term=article_title, retmax=retmax)
    record = Entrez.read(handle)
    handle.close()
    
    # Retrieve the full study articles based on the search results
    id_list = record["IdList"]
    handle = Entrez.efetch(db="pubmed", id=id_list, rettype="medline", retmode="text")
    records = Medline.parse(handle)
    records = list(records)
    handle.close()
    
    # Extract relevant information from the articles and return as JSON or CSV
    articles = []
    for record in records:
        article = {
            "PMID": record["PMID"],
            "Title": record["TI"],
            "Abstract": record.get("AB", ""),
            "Citations": f"https://pubmed.ncbi.nlm.nih.gov/{record['PMID']}/",
        }
        articles.append(article)
    
    # Return the articles as JSON or CSV
    return articles

def search_ctgov(article_title, retmax=5):
    # Perform the ClinicalTrials.gov search using the article title
    api_url = "https://clinicaltrials.gov/api/query/full_studies"
    params = {
        "expr": article_title,
        "min_rnk": 1,
        "max_rnk": retmax,
        "fmt": "json",
    }
    response = requests.get(api_url, params=params)
    data = response.json()
    
    # Extract relevant information from the ctgov results and return as JSON or CSV
    articles = []
    for study in data.get("FullStudiesResponse", {}).get("FullStudies", []):
        article = {
            "PMID": study.get("Study", {}).get("ProtocolSection", {}).get("IdentificationModule", {}).get("NCTId", ""),
            "Title": study.get("Study", {}).get("ProtocolSection", {}).get("IdentificationModule", {}).get("OfficialTitle", ""),
            "Abstract": study.get("Study", {}).get("ProtocolSection", {}).get("DescriptionModule", {}).get("BriefSummary", ""),
            "Citations": f"https://clinicaltrials.gov/ct2/show/{study['Study']['ProtocolSection']['IdentificationModule']['NCTId']}",
        }
        articles.append(article)
    
    # Return the articles as JSON or CSV
    return articles

# Function to display the article details in the main container
def display_articles(articles):
    for article in articles:
        title = article["Title"]
        abstract = article["Abstract"].strip().split(". ", 3)[0] + "..."  # First 3 lines of abstract
        citations_url = article["Citations"]
        st.write(f"**Title:** {title}")
        st.write(f"**Abstract:** {abstract}")
        st.write(f"[Read More]({citations_url})")
        st.write("--------")

@st.cache_resource(ttl="1h")
def configure_retriever(uploaded_file):
    # Read documents

    temp_dir = tempfile.TemporaryDirectory()

    temp_filepath = os.path.join(temp_dir.name, uploaded_file.name)
    with open(temp_filepath, "wb") as f:
        f.write(uploaded_file.getvalue())
    loader = PyPDFLoader(temp_filepath)
    docs = loader.load()

    # Split documents
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1500, chunk_overlap=200)
    splits = text_splitter.split_documents(docs)

    # Create embeddings and store in vectordb
    embeddings = HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")
    vectordb = DocArrayInMemorySearch.from_documents(splits, embeddings)

    # Define retriever
    retriever = vectordb.as_retriever(search_type="mmr", search_kwargs={"k": 2, "fetch_k": 4})

    return retriever


class StreamHandler(BaseCallbackHandler):
    def __init__(self, container: st.delta_generator.DeltaGenerator, initial_text: str = ""):
        self.container = container
        self.text = initial_text

    def on_llm_new_token(self, token: str, **kwargs) -> None:
        self.text += token
        self.container.markdown(self.text)


class PrintRetrievalHandler(BaseCallbackHandler):
    def __init__(self, container):
        self.container = container.expander("Context Retrieval")

    def on_retriever_start(self, query: str, **kwargs):
        self.container.write(f"**Question:** {query}")

    def on_retriever_end(self, documents, **kwargs):
        # self.container.write(documents)
        for idx, doc in enumerate(documents):
            source = os.path.basename(doc.metadata["source"])
            self.container.write(f"**Document {idx} from {source}**")
            self.container.markdown(doc.page_content)

#function to set background image
def set_bg_hack(main_bg):
    '''
    A function to unpack an image from root folder and set as bg.
 
    Returns
    -------
    The background.
    '''
    # set bg name
    main_bg_ext = "png"
        
    st.markdown(
         f"""
         <style>
         .stApp {{
             background: url(data:image/{main_bg_ext};base64,{base64.b64encode(open(main_bg, "rb").read()).decode()});
             background-size: cover
         }}
         </style>
         """,
         unsafe_allow_html=True
     )

def sidebar_bg(side_bg):

   side_bg_ext = 'png'

   st.markdown(
      f"""
      <style>
      [data-testid="stSidebar"] > div:first-child {{
          background: url(data:image/{side_bg_ext};base64,{base64.b64encode(open(side_bg, "rb").read()).decode()});
      }}
      </style>
      """,
      unsafe_allow_html=True,
      )
    
def header_bg(side_bg):

   side_bg_ext = 'png'

   st.markdown(
      f"""
      <style>
      header.css-1avcm0n {{
          background: url(data:image/{side_bg_ext};base64,{base64.b64encode(open(side_bg, "rb").read()).decode()});
      }}
      </style>
      """,
      unsafe_allow_html=True,
      )

def get_sentiment(polarity):
    if polarity < 0.0:
        return 'Negative'
    elif polarity > 0.2:
        return 'Positive'
    else:
        return 'Neutral'
    
def word_frequency(sentence):
    # joins all the sentenses
    #sentence = " ".join(sentence)
    # creates tokens, creates lower class, removes numbers and lemmatizes the words
    new_tokens = word_tokenize(sentence)
    new_tokens = [t.lower() for t in new_tokens]
    new_tokens =[t for t in new_tokens if t not in stopwords.words('english')]
    new_tokens = [t for t in new_tokens if t.isalpha()]
    lemmatizer = WordNetLemmatizer()
    new_tokens =[lemmatizer.lemmatize(t) for t in new_tokens]
    #counts the words, pairs and trigrams
    counted = Counter(new_tokens)
    counted_2= Counter(ngrams(new_tokens,2))
    counted_3= Counter(ngrams(new_tokens,3))
    #creates 3 data frames and returns thems
    word_freq = pd.DataFrame(counted.items(),columns=['word','frequency']).sort_values(by='frequency',ascending=False)
    word_pairs =pd.DataFrame(counted_2.items(),columns=['pairs','frequency']).sort_values(by='frequency',ascending=False)
    trigrams =pd.DataFrame(counted_3.items(),columns=['trigrams','frequency']).sort_values(by='frequency',ascending=False)
    return word_freq,word_pairs,trigrams    
    
#function to read prompt from corresponding text file
def prompt(file):
    with open(file,encoding="utf8") as f:
        return f.read()
    
#function to save a file
def save_uploadedfile(uploaded_file):
     with open(os.path.join(datadir, uploaded_file.name),"wb") as f:
         f.write(uploaded_file.getbuffer())
     return st.success(f"""Saved File:{uploaded_file.name} to directory""")

# def create_vector():
#     documents = SimpleDirectoryReader(saved_path).load_data()
#     index = GPTVectorStoreIndex.from_documents(documents)

#     storage_context = StorageContext.from_defaults()
#     index.storage_context.persist("./vectordatabase")
#     #print ("Done")

# def generate_response(prompt):
#     storage_context = StorageContext.from_defaults(persist_dir="./vectordatabase")
#     index = load_index_from_storage(storage_context)
#     query_engin = index.as_query_engine() 
#     question = prompt
#     response = query_engin.query(question)
#     return str(response)
#     #print ("\n", response)

@st.cache_data
#function to display the PDF of a given file 
def displayPDF(file):
    # Opening file from file path
    with open(file, "rb") as f:
        base64_pdf = base64.b64encode(f.read()).decode('utf-8')

    # Embedding PDF in HTML
    pdf_display = F'<iframe src="data:application/pdf;base64,{base64_pdf}" width="300" height="1100" type="application/pdf"></iframe>'

    # Displaying File
    st.markdown(pdf_display, unsafe_allow_html=True)

# Placeholder function for processing the uploaded documents
def process_documents(NCT, uploaded_file, tense, pls_grade):
    # Implement the document processing logic here

    # Convert tense strings
    tense_mapping = {"on-going": "present", "completed": "past", "upcoming": "future"}
    tense = tense_mapping.get(tense, tense)
    
    summary_replacements = {
        "<Title>": prompt(os.path.join(promptdir, 'title.txt')),
        "<Subtitle>": prompt(os.path.join(promptdir, 'subtitle.txt')),
        "<Key takeaway>": prompt(os.path.join(promptdir, 'key_takeaway.txt')),
        "<Phonetics>": prompt(os.path.join(promptdir, 'phonetics.txt')), 
        "<Introduction>": prompt(os.path.join(promptdir, 'introduction.txt')), 
        "<Intro summary>": prompt(os.path.join(promptdir, 'intro_summary.txt')),
        # "<Inclusion criteria>": "",
        # "<Exclusion crtieria>": "",
        # "<Results>": "",
        "<Aims>": prompt(os.path.join(promptdir, 'aims.txt')),
        "<Conclusions>": prompt(os.path.join(promptdir, 'conclusions.txt')),
        # "<Sponsor>": "",
        # "<More Information>": "",
    }
        
    # Get the text for each section using GPTAPIcall function
    for section_name, summary_prompt in summary_replacements.items():
        
        #prompt for pls grade and tense
        query = f"Strictly following the above instructions and the research document provided, write the content of {section_name} section of the plain language summary in {tense} tense.\
        Do not violate the section-wise instructions provided in any case. The content should be strictly inferred from the research document provided and not any other sources."
        
        st.subheader(f""":red[{section_name[1:-1]} :]""")
        text = llama_vector_index(uploaded_file, prompt(os.path.join(promptdir, f'apls_persona_{pls_grade}_literacy.txt')) + "\n" + summary_prompt + "\n" + query)
        summary_replacements[section_name] = str(text)
        
    ctgov_replacements = {
                    "<Start date>": "Answer the Study Start date in ```MMM-YYYY``` format",
                    "<End date>": "Answer the Study End date in ```MMM-YYYY``` format",
                    "<Participants>": "Total number of Participants in the study including drug arms, placebo arm, soc arm. Give one number answer",
                    "<Arms count>": "Number of arms in the study including the drug arms, placebo arm, soc arm. Give one number answer",
                    "<Disease condition>": "What is the disease condition for which drug is undergoing trials on patients in the study. Give answer as one disease",
                    "<Demographics>": "What are the Demographics of participants in the study",
                    "<treatment arm>": "Number of participants only in the drug arms of the study, do not count the participants from placebo arm or soc arm. Give one number answer",
                    "<control arm>": "Number of participants in the placebo arm or soc arm. Give one number answer",
                    "<Inclusion criteria>": "Inclusion criteria in EligibilityCriteria",
                    "<Exclusion criteria>": "Exclusion criteria in EligibilityCriteria",
                    "<Results>": "list all outcome measure results in bullets interms of outcome measure type, outcome measure title, outcome measure description, outcome measure value",
                    # "<clinical trials gov link>": "https://clinicaltrials.gov/ct2/show/NCT03036813",
                    # "<Summary date>": datetime.now().strftime('%d-%b-%Y'),
                    "<Sponsor>": "Lead Sponsor Name",
                   }
    
    for section_name, ctgov_prompt in ctgov_replacements.items():
        
        st.subheader(f""":red[{section_name[1:-1]} :]""")
        text = CTGovAPIcall(NCT, ctgov_prompt)
        if section_name=="<Participants>":
            text = re.findall(r'\d+', text)
        ctgov_replacements[section_name] = str(text)
    
    
    replacements = {**summary_replacements, 
                    **ctgov_replacements, 
                    "<Study number>": f"{NCT}",
                    "<clinical trials gov link>": f"https://clinicaltrials.gov/ct2/show/{NCT}",
                    "<Summary date>": datetime.now().strftime('%d-%b-%Y'),
                   }
    
    return replacements

def CTGovAPIcall(NCT, query):
    file_format = '&fmt=JSON'
    expr = NCT #'A+Phase+3+Randomized+Trial+of+Voxelotor+in+Sickle+Cell+Disease' #or give NCT number here NCT03036813
    ctgov = 'https://classic.clinicaltrials.gov/api/query/full_studies?expr='

    your_url = (ctgov + expr + file_format)

    with urllib.request.urlopen(your_url) as url:
        ini_dict = json.loads(url.read().decode())
        
    json_spec = JsonSpec(dict_=ini_dict["FullStudiesResponse"]["FullStudies"][0]["Study"], max_value_length=31000)
    json_toolkit = JsonToolkit(spec=json_spec)
    
    # chat_box = st.empty()
    # stream_handler = StreamHandler(chat_box, display_method='write')
    
    json_agent_executor = create_json_agent(
        llm=ChatOpenAI(temperature=0, model_name="gpt-4-32k", streaming=True, callbacks=[StreamlitCallbackHandler(st.container())],), toolkit=json_toolkit, verbose=True
    )
    resp = json_agent_executor.run(query)
    st.write(resp)
    return resp
    
# Placeholder function for postprocessing into PPT template
def postprocess_to_ppt(replacements, selected_template):
    # Implement the postprocessing logic here
    # For demonstration purposes, we'll load a presentation object and copy the text from replacements dictionary
    
    #rootdir = os.path.realpath('./')
    
    #selected_template = "PLS_PPT_Template"
    selected_template = selected_template[:-4]
    ppt_file = f"{selected_template}.pptx"
    prs = Presentation(os.path.join(rootdir, ppt_file))

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        for placeholder, new_text in replacements.items():
                            if run.text == placeholder:
                                # Preserve formatting of the first run in the paragraph
                                first_run = paragraph.runs[0]
                                font_size = first_run.font.size
                                font_name = first_run.font.name
                                font_bold = first_run.font.bold
                                font_italic = first_run.font.italic

                                # Check if font color is explicitly defined
                                if first_run.font.color.type == "rgb":
                                    font_color = first_run.font.color.rgb
                                else:
                                    font_color = None

                                # Replace text while preserving formatting
                                run.text = new_text

                                # Apply preserved formatting to the entire paragraph
                                for run in paragraph.runs:
                                    run.font.size = font_size
                                    run.font.name = font_name
                                    run.font.bold = font_bold
                                    run.font.italic = font_italic
                                    if font_color:
                                        run.font.color.rgb = font_color

    # Return the modified presentation object
    return prs


# Placeholder function for postprocessing into DOC template
def postprocess_to_doc(replacements):
    
    para_variable_list = ["<Subtitle>", "<Key takeaway>", "<Phonetics>", "<Introduction>", "<Intro summary>", "<Demographics>", "<Inclusion criteria>", "<Exclusion criteria>", "<Results>", "<Aims>", "<Conclusions>"]
    table_variable_list = ["<Study number>", "<Start date>", "<End date>", "<Participants>", "<Arms count>", "<treatment arm>", "<control arm>", "<Sponsor>", "<Summary date>", "<clinical trials gov link>"]
    # Create a new document
    document = Document()

    # Set the font size of the document
    style = document.styles['Normal']
    font = style.font
    font.size = Pt(11)

    # Set the title
    title = replacements.get("<Title>")
    if title:
        document.add_heading(title, level=1).bold = True

    # Add paragraphs for para_variable_list with the same header formatting
    for variable in para_variable_list:
        value = replacements.get(variable)
        if value:
            p = document.add_paragraph(style='Heading 1')
            p.text = variable[1:-1]
            p.bold = True
            document.add_paragraph(value)

    # Add the table for table_variable_list
    table_replacements = {variable: replacements.get(variable) for variable in table_variable_list}
    if table_replacements:
        table_heading = "Additional Information"
        document.add_heading(table_heading, level=1)

        # Create the table
        table = document.add_table(rows=1, cols=2)
        table.style = 'Table Grid'
        
        # Set table column widths
        table.autofit = False
        table.columns[0].width = Pt(200)
        table.columns[1].width = Pt(300)

        # Add table headers
        table_header_cells = table.rows[0].cells
        table_header_cells[0].text = "Variable"
        table_header_cells[1].text = "Value"
        for cell in table_header_cells:
            cell.paragraphs[0].alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
            cell.paragraphs[0].bold = True

        # Add table rows
        for variable, value in table_replacements.items():
            row_cells = table.add_row().cells
            row_cells[0].text = variable[1:-1]
            row_cells[1].text = value

    return document
    
    
def main():
    
    #Page icons n tab name on browser tab
    #img = Image.open(os.path.join(rootdir, 'pfizer.png'))
    st.set_page_config(page_title = 'MAIA', page_icon = ":robot_face:", layout="wide")
    
    #to hide the hamburger running on top right and footer of streamlit
    # hide_default_format = """
    #    <style>
    #    #MainMenu {visibility: hidden; }
    #    footer {visibility: hidden;}
    #    </style>
    #    """
    # st.markdown(hide_default_format, unsafe_allow_html=True)
    
    names = ["admin","shakti"]
    usernames = ["adm", "shrp"]
    passwords = ["abc123", "def456"]

    credentials = {"usernames":{}}
    hashed_passwords = stauth.Hasher(passwords).generate()
    
    for uname, name, pwd in zip(usernames, names, hashed_passwords):
        user_dict = {"name": name, "password": pwd}
        credentials["usernames"].update({uname: user_dict})

    
    #add a cookie which will be stored on client browser to save credentials till 30days
    authenticator = stauth.Authenticate(credentials, "pls_generator", "abcdef", cookie_expiry_days = 30)

    #u can locate the authenticator in the main body or the sidebar
    name, authentication_status, username = authenticator.login("Login", "main")
    
    if st.session_state["authentication_status"] == False:
        st.error("Username/password is incorrect")
        
    if st.session_state["authentication_status"] == None:
        st.warning("Please enter your username and password")
        
    if st.session_state["authentication_status"]:
        
        #logout button on main container
        authenticator.logout('Logout', 'main')
        st.subheader(f'Welcome *{st.session_state["name"]}*')
        #st.session_state.openai_api_key  = st.text_input("Enter your OpenAI API Key", '',type="password")
        
        #set bg image cover
        #set_bg_hack(os.path.join(rootdir, 'iqvia-dark-blue.png'))
        sidebar_bg(os.path.join(rootdir, 'iqvia-blue.png'))
        #header_bg(os.path.join(rootdir, 'iqvia-dark-blue.png'))

        #setting banner image
        #st.image(Image.open(os.path.join(rootdir, 'Pfizer-AI.jpg')))
        
        selected_tab = option_menu(
            menu_title=None,  # required
            options=["PLS Generator", "Search your Document", "Word Analytics", "Converse with your Documents", "Search and Chat with PubMed/CTGov"],  # required
            icons=["house", "book", "envelope", "book", "envelope"],  # optional
            menu_icon="cast",  # optional
            default_index=0,  # optional
            orientation="horizontal",
            # styles={
            #     "container": {"padding": "0!important"},
            #     "icon": {"color": "orange", "font-size": "25px"},
            #     "nav-link": {
            #         "font-size": "25px",
            #         "text-align": "left",
            #         "margin": "0px",
            #         "--hover-color": "#eee",
            #     },
            #     "nav-link-selected": {"background-color": "green"},
            # },
        )
        
        #setting input components on sidebar
        with st.sidebar:

            st.image(Image.open(os.path.join(rootdir, 'iqvia-logo.png')))
            #setting title
            st.markdown("""<h3 style='text-align: center'>*MAIA - Medical Affairs Intelligence Assistant*</h3>""", unsafe_allow_html=True)

            # Step 1: Document Upload
            st.subheader("Step 1: Upload Clinical trial document")
            uploaded_file = st.file_uploader("Upload document", accept_multiple_files=False, type=["pdf"])
            
            NCT = st.text_input("Enter the NCT number:", "NCT", key="NCT")
            
            # Step 2: User Inputs
            st.subheader("Step 2: Define the tone and Grade of PLS")
            # Set default values for radio button and slider
            default_tense = "Completed"
            default_pls_grade = "Low"

            # Radio button for tense selection
            tense = st.radio("Current status of the study for writing tense", options=["On-going", "Completed", "Upcoming"], key="tense", index=["On-going", "Completed", "Upcoming"].index(default_tense), horizontal=True)
            #st.write('<style>div.row-widget.stRadio > div{flex-direction:row;}</style>', unsafe_allow_html=True)

            # Slider for PLS grade selection
            #pls_grade = st.slider("Health Literacy Grade Reading level", min_value=0, max_value=10, step=5, key="pls_grade", value=default_pls_grade)
            pls_grade = st.select_slider("Health Literacy Grade of audience", options=["Low", "High"], key="pls_grade", value = default_pls_grade)

            st.session_state.process_button = False
            process_button = st.button("Process Documents")
            st.session_state.process_button = process_button
            st.session_state.uploaded_file = uploaded_file
            st.session_state.selected_tab = selected_tab
        
        #if st.session_state.process_button and st.session_state.uploaded_file:
        #if process_button and uploaded_file:
            
            # Retrieve user inputs if you haven't initialized them to any variable, then retrieve from streamlit session state
            # st.session_state.tense = tense
            # st.session_state.pls_grade = pls_grade
            
        if st.session_state.selected_tab == "Search your Document":
            st.subheader("Ask your PDF 💬")
            # show user input
            user_question = st.text_input("Ask a question about your PDF:", placeholder="Number of participants? ", disabled=not uploaded_file,)
            
            if st.session_state.uploaded_file:              
                
                # extract the text
                if uploaded_file is not None:
                  pdf_reader = PdfReader(uploaded_file)
                  text = ""
                  for page in pdf_reader.pages:
                    text += page.extract_text()

                  # split into chunks
                  text_splitter = CharacterTextSplitter(
                    separator="\n",
                    chunk_size=1000,
                    chunk_overlap=200,
                    length_function=len
                  )
                  chunks = text_splitter.split_text(text)
                    
                  # create embeddings
                  store_name = uploaded_file.name[:-4]
                  if os.path.exists(os.path.join(datadir, f"{store_name}.pkl")):
                    with open(os.path.join(datadir, f"{store_name}.pkl"), "rb") as f:
                        knowledge_base = pickle.load(f)
                        st.write('Embeddings loaded from the Disk:')
                  else:
                    embeddings = OpenAIEmbeddings()
                    knowledge_base = FAISS.from_texts(chunks, embeddings)
                    with open(os.path.join(datadir, f"{store_name}.pkl"), "wb") as f:
                        pickle.dump(knowledge_base, f)
                        st.write('Embeddings newly created')

                  if user_question:
                    docs = knowledge_base.similarity_search(user_question, k=3)

                    # chat_box = st.empty()
                    # stream_handler = StreamHandler(chat_box, display_method='write')
        
                    llm = ChatOpenAI(temperature=0, callbacks=[StreamlitCallbackHandler(st.container())], streaming=True)
                    chain = load_qa_chain(llm, chain_type="stuff")
                    
                    #get_openai_callback() gives the cost on console
                    # with get_openai_callback() as cb:
                    #   response = chain.run(input_documents=docs, question=user_question)
                    #   print(cb)
                    response = chain.run(input_documents=docs, question=user_question)
                    st.write(response)    
        
        if st.session_state.selected_tab == "Converse with your Documents":
#             # Ensure the directory exists
#             if not os.path.exists(saved_path):
#                 os.makedirs(saved_path)
        
#             if uploaded_file is not None:
#                 # To read file as bytes:
#                 bytes_data = uploaded_file.getvalue()

#                 # Save the uploaded file to the 'data' directory
#                 with open(os.path.join(saved_path, uploaded_file.name), 'wb') as out_file:
#                     out_file.write(bytes_data)

#                 st.success('PDF file saved in data directory')
#                 create_vector()
#                 #remove_all_files(saved_path)
#                 st.success('Vector created')

#             # Initialise session state variables
#             if 'generated' not in st.session_state:
#                 st.session_state['generated'] = []
#             if 'past' not in st.session_state:
#                 st.session_state['past'] = []
#             if 'messages' not in st.session_state:
#                 st.session_state['messages'] = [
#                     {"role": "system", "content": "You are a helpful assistant."}
#                 ]

#             response_container = st.container()
#             # container for text box
#             container = st.container()

#             with container:
#                 with st.form(key='my_form', clear_on_submit=True):
#                     user_input = st.text_area("You:", key='input', height=50)
#                     submit_button = st.form_submit_button(label='Send')

#                 if submit_button and user_input:
#                     output = generate_response(user_input)
#                     st.session_state['past'].append(user_input)
#                     st.session_state['generated'].append(output)
#                     #st.session_state['model_name'].append(model_name)

#             if st.session_state['generated']:
#                 with response_container:
#                     for i in range(len(st.session_state['generated'])):
#                         message(st.session_state["past"][i], is_user=True, key=str(i) + '_user1')
#                         message(st.session_state["generated"][i], key=str(i))
            
            if not uploaded_file:
                st.caption("Please upload research PDF documents to continue.")
                st.stop()
    
            retriever = configure_retriever(uploaded_file)

            # Setup memory for contextual conversation
            msgs = StreamlitChatMessageHistory()
            memory = ConversationBufferMemory(memory_key="chat_history", chat_memory=msgs, return_messages=True)

            # Setup LLM and QA chain
            llm = ChatOpenAI(
                model_name="gpt-3.5-turbo", openai_api_key=st.secrets["openai_api_key"], temperature=0, streaming=True
            )
            qa_chain = ConversationalRetrievalChain.from_llm(
                llm, retriever=retriever, memory=memory, verbose=True
            )

            if len(msgs.messages) == 0 or st.sidebar.button("Clear message history"):
                msgs.clear()
                msgs.add_ai_message("How can I help you?")

            avatars = {"human": "user", "ai": "assistant"}
            for msg in msgs.messages:
                st.chat_message(avatars[msg.type]).write(msg.content)

            if user_query := st.chat_input(placeholder="Ask me anything!"):
                st.chat_message("user").write(user_query)

                with st.chat_message("assistant"):
                    retrieval_handler = PrintRetrievalHandler(st.container())
                    stream_handler = StreamHandler(st.empty())
                    response = qa_chain.run(user_query, callbacks=[retrieval_handler, stream_handler])
                
    
    
        if st.session_state.selected_tab == "Word Analytics":    
            #st.subheader(f"You have selected {selected_tab}")
            
            # extract the text
            if uploaded_file is not None:
                pdf_reader = PdfReader(uploaded_file)
                text = ""
                for page in pdf_reader.pages:
                    text += page.extract_text()
                
                #Remove un-important words:
                stop_words = set(stopwords.words('english'))
                query_words={'participants', 'Participants' }
                stop_words.update(query_words)
                for word in query_words:
                    text = text.replace(word, '')
                    
                # Create and generate a word cloud image:
                wordcloud = WordCloud(stopwords=stop_words, collocations=False, max_font_size=55, max_words=25, background_color="black").generate(text)

                # Display the generated image:
                plt.figure(figsize=(10,12))
                plt.imshow(wordcloud, interpolation='bilinear')
                plt.axis("off")
                plt.show()
                st.set_option('deprecation.showPyplotGlobalUse', False)
                st.pyplot()
                
                col1, col2 = st.columns([0.3,0.7])
                
                with col1:
                    #overall doc sentiment
                    analyzer=SentimentIntensityAnalyzer()
                    polarity = analyzer.polarity_scores(text)['compound']
                    st.subheader(f"\nOverall Research document sentiment is: {get_sentiment(polarity)}", )
            
                    # splitting modified single string into list of strings using groupby() function
                    grouped_strings = ["".join(g) for k, g in itertools.groupby(text, lambda x: x == " ") if not k]
                    
                    #word-wise sentiment
                    df = pd.DataFrame()
                    df['polarity']=[analyzer.polarity_scores(text)['compound'] for text in grouped_strings]
                    df['sentiment']=df.polarity.apply(get_sentiment)
                    plt.figure(figsize=(3,3))
                    df.sentiment.value_counts().plot.bar()
                    st.pyplot()
                    
                    #Tokenization
                    tokens = nltk.word_tokenize(text)

                    #compute freq distribution
                    freq_dist = FreqDist(tokens)

                    #plot the freq distribution
                    freq_dist.plot(50, cumulative=True)

                    #set labels and title
                    # plt.xlabel('Words')
                    # plt.ylabel('Frequency')
                    plt.title('Frequency distribution of words')
                    st.pyplot()
                    
                with col2:
                    st.subheader("\nWord frequency of the words in the research doc is: ")
                    
                    data2,data3,data4 = word_frequency(text)
                    fig, axes = plt.subplots(3,1,figsize=(8,20))
                    sns.barplot(ax=axes[0],x='frequency',y='word',data=data2.head(30))
                    sns.barplot(ax=axes[1],x='frequency',y='pairs',data=data3.head(30))
                    sns.barplot(ax=axes[2],x='frequency',y='trigrams',data=data4.head(30))
                    st.pyplot(fig)
                
        if st.session_state.selected_tab == "PLS Generator":
            if st.session_state.process_button and st.session_state.uploaded_file and st.session_state.NCT!='NCT':
                col1, col2 = st.columns([0.2,0.8], gap="large")
                with col1:
                    input_file = save_uploadedfile(uploaded_file)
                    pdf_file = os.path.join(datadir, uploaded_file.name) #rootdir + "/" + uploaded_file.name
                    pdf_view = displayPDF(pdf_file)
                with col2:
                    with st.spinner(text='Processing research document you gave on the left to generate Plain Language Summary for you...⏳'):

                        # Progress bar
                        #progress_bar_method(50) or st.progress(0, "text")
                        
                        # Call the processing function on the uploaded documents with user inputs
                        replacements = process_documents(NCT, pdf_file, tense, pls_grade)
                        st.success("Processed Output to be filled up in the preferred PLS template")

                        #Display processed output
                        #st.write(replacements)
                        st.snow()
                        st.balloons()
                        
                # Store the replacements dictionary in session state
                st.session_state.replacements = replacements

            # Step 3: PPT Template Selection and Download
            st.subheader("Step 3: Select PLS Template and Download")
            
            default_format = "PPT format"
            st.session_state.select_format = pills("Select PPT or Word format", ["PPT format", "Word format"], ["🎈", "🌈"], index=["PPT format", "Word format"].index(default_format))
            
            if st.session_state.select_format == "PPT format":
                # Add radio buttons for template selection here    
                default_template = "Blue_PLS_Template.png"
                selected_template = image_select(
                    label="Select PPT Template",
                    images=[
                        os.path.join(rootdir, 'Pfizer_Blue_PLS_Template.png'),
                        os.path.join(rootdir, 'Pfizer_Red_PLS_Template.png'),
                        os.path.join(rootdir, 'Pfizer_Long_PLS_Template.png'),
                    ],
                    captions=["Blue_PLS_Template", "Red_PLS_Template", "Long_PLS_Template"],
                    index=["Blue_PLS_Template.png", "Red_PLS_Template.png", "Long_PLS_Template.png"].index(default_template),
                    use_container_width = False,
                )
                #selected_template = st.radio("Select PPT Template", options=["Blue_PLS_Template", "Red_PLS_Template", "Long_PLS_Template"], index=["Blue_PLS_Template", "Red_PLS_Template", "Long_PLS_Template"].index(default_template), horizontal=True)
                #selected_template = pills("", ["Pfizer_Blue_PLS_Template", "Pfizer_Red_PLS_Template", "Pfizer_Long_PLS_Template"], ["🍀", "🎈", "🌈"])
            
            if st.session_state.select_format == "Word format":
                default_template = "Word_PLS_Template"
                selected_template = "Blue_PLS_Template.png"
                image_select(
                    label="Select Word Template",
                    images=[
                        os.path.join(rootdir, 'Pfizer_Word_PLS_Template.png'),
                    ],
                    captions=["Word_PLS_Template"],
                    index=["Word_PLS_Template"].index(default_template),
                    use_container_width = False,
                )
            
            generate_ppt_button = st.button("Generate PLS")

            if generate_ppt_button:
                # Retrieve the replacements dictionary from session state
                replacements = st.session_state.replacements
                st.session_state.process_button = False

                if replacements:
                    with st.spinner('Generating PLS slides for you...⏳'):
                        # Call the postprocessing function to generate PPT content
                        ppt_content = postprocess_to_ppt(replacements, selected_template)

                        doc_content = postprocess_to_doc(replacements)

                        # Display the PPT content using st.markdown or st.write
                        #st.markdown(ppt_content, unsafe_allow_html=True)
                        st.markdown(list(replacements.keys()))

                        # Store the PPT content in session state
                        st.session_state.ppt_content = ppt_content
                        st.session_state.doc_content = doc_content

                         # Step 4: PPT Download
                        if "ppt_content" and "doc_content" in st.session_state:
                            ppt_content = st.session_state.ppt_content
                            doc_content = st.session_state.doc_content

                            st.session_state.replacements = replacements
                            st.session_state.process_button = False

                            # Save the modified presentation object to a temporary file
                            #ppt_output_file = f"PLS_{replacements['<Title>']}_{datetime.now().strftime('%Y%m%d%H%M%S')}.pptx"                    
                            ppt_output_file = "PLS_PPT.pptx"
                            #ppt_content.save(ppt_output_file)

                            # save presentation as binary output
                            binary_output = BytesIO()
                            ppt_content.save(binary_output)

                            binary_output_doc = BytesIO()
                            doc_content.save(binary_output_doc)

                            # display success message and download button
                            st.success(':tada: The PLS template has been filled with above sections in ' + selected_template)

                            # Provide the download link for the generated PPT and DOC
                            if st.session_state.select_format == "PPT format":
                                st.download_button("Download PLS PPT", data=binary_output.getvalue(), file_name=ppt_output_file, mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
                            if st.session_state.select_format == "Word format":
                                st.download_button("Download PLS Doc", data=binary_output_doc.getvalue(), file_name="PLS_DOC.docx", mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document")

                                
        if st.session_state.selected_tab == "Search and Chat with PubMed/CTGov":
            st.subheader("Search PubMed and ClinicalTrials.gov")
            query = st.text_input("Enter your query:")
            search_option = st.radio("Search Option", ["Search PubMed", "Search ClinicalTrials.gov", "Search Both"], horizontal=True)

            col1, col2 = st.columns(2)  # Split the main container into two columns
            if query:
                if search_option == "Search PubMed":
                    pubmed_articles = search_pubmed(query)[:5]
                    st.session_state.articles = pubmed_articles
                    with col1:
                        st.write("### PubMed Results")
                        display_articles(pubmed_articles)
                    # Create a DataFrame in the second column with PubMed articles
                    with col2:
                        df_pubmed = pd.DataFrame(pubmed_articles)
                        st.write("### Converse with retrieved PubMed Articles")
                        st.write(df_pubmed)

                elif search_option == "Search ClinicalTrials.gov":
                    ctgov_articles = search_ctgov(query)[:5]
                    st.session_state.articles = ctgov_articles
                    with col1:
                        st.write("### ClinicalTrials.gov Results")
                        display_articles(ctgov_articles)
                    # Create a DataFrame in the second column with ClinicalTrials.gov articles
                    with col2:
                        df_ctgov = pd.DataFrame(ctgov_articles)
                        st.write("### Converse with retrieved CTGov Articles")
                        st.write(df_ctgov)

                else:  # Search Both
                    pubmed_articles = search_pubmed(query)[:5]
                    ctgov_articles = search_ctgov(query)[:5]
                    st.session_state.articles = pubmed_articles + ctgov_articles
                    with col1:
                        st.write("### PubMed Results")
                        display_articles(pubmed_articles)
                    with col1:  # Use the same column for ClinicalTrials.gov Results
                        st.write("### ClinicalTrials.gov Results")
                        display_articles(ctgov_articles)
                    # Concatenate both sets of articles and create a DataFrame in the second column
                    with col2:
                        combined_articles = pubmed_articles + ctgov_articles
                        df_combined = pd.DataFrame(combined_articles)
                        st.write("### Converse with retrieved PubMed/CTGov Articles")
                        st.write(df_combined)
            # Placeholder for chatbot implementation in the second column
            with col2:
                st.write("Chatbot - Ask questions from only these Pubmed/CTGov articles")
                if 'articles' in st.session_state:      
                    df = pd.DataFrame(st.session_state.articles)
                    #user_prompt = st.text_area(label="prompt:",placeholder="Number of patients..",)
                    #if st.button("Generate"):
                ########################################################LangChain CSV Agent (with Pandas)

                    # langchain_pandas_agent = create_pandas_dataframe_agent(
                    #     ChatOpenAI(temperature=0, model="gpt-4-32k", streaming=True, ),
                    #     df,
                    #     verbose=True,
                    #     agent_type=AgentType.OPENAI_FUNCTIONS,
                    # )

                    #st.write("Langchain pandas agent: ", langchain_pandas_agent.run(user_prompt))
                #####################################for only QnA, not chatbot########use code in this brackets
#                     # Initialise session state variables
#                     if 'generated1' not in st.session_state:
#                         st.session_state['generated1'] = []
#                     if 'past1' not in st.session_state:
#                         st.session_state['past1'] = []
#                     if 'messages1' not in st.session_state:
#                         st.session_state['messages1'] = [
#                             {"role": "system", "content": "You are a helpful assistant."}
#                         ]


#                     # container for chat history
#                     response_container = st.container()

#                     # container for text box
#                     input_container = st.container()

#                     with input_container:
#                         # Create a form for user input
#                         with st.form(key='my_form', clear_on_submit=True):
#                             user_input = st.text_area("You:", key='input', height=100)
#                             submit_button = st.form_submit_button(label='Send')

#                         if submit_button and user_input:
#                             # If user submits input, generate response and store input and response in session state variables
#                             try:
#                                 query_response = generate_response1(user_input, df)
#                                 st.session_state['past1'].append(user_input)
#                                 st.session_state['generated1'].append(query_response)
#                             except Exception as e:
#                                 st.error("An error occurred: {}".format(e))

#                     if st.session_state['generated1']:
#                         # Display chat history in a container
#                         with response_container:
#                             for i in range(len(st.session_state['generated1'])):
#                                 message(st.session_state["past1"][i], is_user=True, key=str(i) + '_user')
#                                 message(st.session_state["generated1"][i], key=str(i))
                        
        
                    if "messages" not in st.session_state or st.sidebar.button("Clear conversation history"):
                        st.session_state["messages"] = [{"role": "assistant", "content": "How can I help you?"}]

                    for msg in st.session_state.messages:
                        st.chat_message(msg["role"]).write(msg["content"])

                    if prompt := st.text_input("You:",placeholder="What is this data about?"):
                        st.session_state.messages.append({"role": "user", "content": prompt})
                        st.chat_message("user").write(prompt)

                        llm = ChatOpenAI(
                            temperature=0, model="gpt-3.5-turbo-0613", openai_api_key=st.secrets["openai_api_key"], streaming=True
                        )

                        pandas_df_agent = create_pandas_dataframe_agent(
                            llm,
                            df,
                            verbose=True,
                            agent_type=AgentType.OPENAI_FUNCTIONS,
                            handle_parsing_errors=True,
                        )

                        with st.chat_message("assistant"):
                            st_cb = StreamlitCallbackHandler(st.container(), expand_new_thoughts=False)
                            response = pandas_df_agent.run(st.session_state.messages, callbacks=[st_cb])
                            st.session_state.messages.append({"role": "assistant", "content": response})
                            st.write(response)
                            
                            
                        # Add a download button for the chat conversation
                        #if st.button("Download Chat Conversation"):
                            #download_chat_conversation(st.session_state['past'], st.session_state['generated'])
                                
if __name__ == "__main__":
    main()
