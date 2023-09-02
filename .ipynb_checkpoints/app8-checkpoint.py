#enhancing - current working version
import os
import urllib
import base64
import json
import streamlit as st
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from docx import Document
from docx.shared import Pt
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from datetime import datetime  # Import the 'datetime' class from the 'datetime' module
import time
from io import BytesIO
from PIL import Image
from shakti_stream_index import llama_vector_index
from streamlit_pills import pills
import streamlit_authenticator as stauth
from streamlit_option_menu import option_menu
from PyPDF2 import PdfReader
from langchain.text_splitter import CharacterTextSplitter
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain.vectorstores import FAISS
from langchain.chains.question_answering import load_qa_chain
from langchain.llms import OpenAI
from langchain.callbacks import get_openai_callback


rootdir = "/home/cdsw/experimentation_project1/PLS_project"
datadir = "/home/cdsw/experimentation_project1/PLS_project/data"

#function to set background image
def set_bg_hack(main_bg):
    '''
    A function to unpack an image from root folder and set as bg.
 
    Returns
    -------
    The background.
    '''
    # set bg name
    main_bg_ext = "jpg"
        
    st.markdown(
         f"""
         <style>
         .stApp {{
             background: url(data:image/{main_bg_ext};base64,{base64.b64encode(open(main_bg, "rb").read()).decode()});
             background-size: cover
         }}
         </style>
         """,
         unsafe_allow_html=True
     )

def sidebar_bg(side_bg):

   side_bg_ext = 'png'

   st.markdown(
      f"""
      <style>
      [data-testid="stSidebar"] > div:first-child {{
          background: url(data:image/{side_bg_ext};base64,{base64.b64encode(open(side_bg, "rb").read()).decode()});
      }}
      </style>
      """,
      unsafe_allow_html=True,
      )
    
def header_bg(side_bg):

   side_bg_ext = 'png'

   st.markdown(
      f"""
      <style>
      header.css-1avcm0n {{
          background: url(data:image/{side_bg_ext};base64,{base64.b64encode(open(side_bg, "rb").read()).decode()});
      }}
      </style>
      """,
      unsafe_allow_html=True,
      )
    
#function to read prompt from corresponding text file
def prompt(file):
    with open(file) as f:
        return f.read()
    
#function to save a file
def save_uploadedfile(uploaded_file):
     with open(os.path.join(datadir, uploaded_file.name),"wb") as f:
         f.write(uploaded_file.getbuffer())
     return st.success(f"""Saved File:{uploaded_file.name} to directory""")

@st.cache_data
#function to display the PDF of a given file 
def displayPDF(file):
    # Opening file from file path
    with open(file, "rb") as f:
        base64_pdf = base64.b64encode(f.read()).decode('utf-8')

    # Embedding PDF in HTML
    pdf_display = F'<iframe src="data:application/pdf;base64,{base64_pdf}" width="300" height="1100" type="application/pdf"></iframe>'

    # Displaying File
    st.markdown(pdf_display, unsafe_allow_html=True)

# Placeholder function for processing the uploaded documents
def process_documents(uploaded_file, tense, pls_grade):
    # Implement the document processing logic here

    '''
        replacements = {
            "<Title>": GPTAPIcall(uploaded_file, "'Title'", tense, pls_grade),
        #     "<Subtitle>": GPTAPIcall(uploaded_file, "subtitle", tense, pls_grade),
            "<Introduction>": GPTAPIcall(uploaded_file, "'Introduction'", tense, pls_grade),
        #     "<Phonetics>": GPTAPIcall(uploaded_file, "phonetics", tense, pls_grade),
        #    "<Key takeaway>": GPTAPIcall(uploaded_file, "'Key Takeaway'", tense, pls_grade),
        #    "<Results>": GPTAPIcall(uploaded_file, "'Results'", tense, pls_grade)
        #     "<Intro summary>": GPTAPIcall(uploaded_file, "conclusion", tense, pls_grade)
        }
    '''
    replacements = {
        "<Title>": "",
        "<Subtitle>": "",
        "<Key takeaway>": "",
        "<Phonetics>": "",
        "<Introduction>": "",
        "<Intro summary>": "",
        "<Inclusion criteria>": "",
        "<Exclusion crtieria>": "",
        "<Results>": "",
        "<Aims>": "",
        "<Conclusions>": "",
        "<Sponsor>": "",
        "<More Information>": "",
    }
        
    # Get the text for each section using GPTAPIcall function
    for section_name in replacements:
        
        st.subheader(f""":red[{section_name[1:-1]} :]""")
        text = GPTAPIcall(uploaded_file, section_name, tense, pls_grade)
        replacements[section_name] = str(text)
        
    replacements = {**replacements, 
                    "<Participants>": "274",
                    "<Disease condition>": "Sickle cell disease",
                    "<Demographics>": "Aged 12 to 65 years",
                    "<treatment arm>": "182",
                    "<control arm>": "92",
                    "<Study number>": "NCT03036813",
                    "<Start date>": "April 2018",
                    "<End date>": "April 2021",
                    "<clinical trials gov link>": "https://clinicaltrials.gov/ct2/show/NCT03036813",
                    "<Summary date>": datetime.now().strftime('%d-%b-%Y'),
                   }
    
    return replacements

# Placeholder function for GPT API call
def GPTAPIcall(uploaded_file, key, tense, pls_grade):
    # Placeholder logic to generate values based on the key, tense, and PLS grade
    # Replace this with your actual GPT API call or any other processing logic
    # Convert tense strings
    tense_mapping = {"on-going": "present", "completed": "past", "upcoming": "future"}
    tense = tense_mapping.get(tense, tense)

    # Placeholder value for the key with the tense and PLS grade
    #value = f"Placeholder value for {key} (Tense: {tense}, PLS Grade: {pls_grade})"
    query = f"Strictly following the above instructions and the clinical trial document provided, write the content of {key} section of the APLS in {tense}, comprehendable by a {pls_grade} health literacy grade person. Do not violate the section-wise instructions provided in any case. The content should be strictly inferred from the clinical trial document provided only and not any other sources."
    
    return llama_vector_index(uploaded_file, prompt(os.path.join(rootdir, 'apls_persona_2606.txt')) + "\n" + query)

# Placeholder function for postprocessing into PPT template
def postprocess_to_ppt(replacements, selected_template):
    # Implement the postprocessing logic here
    # For demonstration purposes, we'll load a presentation object and copy the text from replacements dictionary
    
    #rootdir = os.path.realpath('./')
    
    #selected_template = "PLS_PPT_Template"
    ppt_file = f"{selected_template}.pptx"
    prs = Presentation(os.path.join(rootdir, ppt_file))

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        for placeholder, new_text in replacements.items():
                            if run.text == placeholder:
                                # Preserve formatting of the first run in the paragraph
                                first_run = paragraph.runs[0]
                                font_size = first_run.font.size
                                font_name = first_run.font.name
                                font_bold = first_run.font.bold
                                font_italic = first_run.font.italic

                                # Check if font color is explicitly defined
                                if first_run.font.color.type == "rgb":
                                    font_color = first_run.font.color.rgb
                                else:
                                    font_color = None

                                # Replace text while preserving formatting
                                run.text = new_text

                                # Apply preserved formatting to the entire paragraph
                                for run in paragraph.runs:
                                    run.font.size = font_size
                                    run.font.name = font_name
                                    run.font.bold = font_bold
                                    run.font.italic = font_italic
                                    if font_color:
                                        run.font.color.rgb = font_color

    # Return the modified presentation object
    return prs


# Placeholder function for postprocessing into DOC template
def postprocess_to_doc(replacements):
    # Create a new document
    document = Document()

    # Set the font size of the document
    style = document.styles['Normal']
    font = style.font
    font.size = Pt(11)

    # Set the title
    title = replacements.get("Title")
    if title:
        document.add_heading(title, level=1).bold = True

    # Add sections and paragraphs
    section_count = 0
    for key, value in replacements.items():
        if key != "Title":
            section_count += 1
            if section_count <= 12:
                document.add_heading(key, level=1)
                if value:
                    document.add_paragraph(value)

    # Add the table
    table_replacements = {k: v for k, v in replacements.items() if k != "Title" and k not in list(replacements.keys())[1:13]}
    if table_replacements:
        table_heading = "Additional Information"
        document.add_heading(table_heading, level=1)

        # Create the table
        table = document.add_table(rows=1, cols=2)
        table.style = 'Table Grid'
        
        # Set table column widths
        table.autofit = False
        table.columns[0].width = Pt(200)
        table.columns[1].width = Pt(300)

        # Add table headers
        table_header_cells = table.rows[0].cells
        table_header_cells[0].text = "Variable"
        table_header_cells[1].text = "Value"
        for cell in table_header_cells:
            cell.paragraphs[0].alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
            cell.paragraphs[0].bold = True

        # Add table rows
        for key, value in table_replacements.items():
            row_cells = table.add_row().cells
            row_cells[0].text = key
            row_cells[1].text = value

    return document
    
    
    

def main():
    
    #Page icons n tab name on browser tab
    #img = Image.open(os.path.join(rootdir, 'pfizer.png'))
    st.set_page_config(page_title = 'MAIA', page_icon = ":robot_face:", layout="wide")
    
    names = ["admin","shakti"]
    usernames = ["adm", "shrp"]
    passwords = ["abc123", "def456"]

    credentials = {"usernames":{}}
    hashed_passwords = stauth.Hasher(passwords).generate()
    
    for uname, name, pwd in zip(usernames, names, hashed_passwords):
        user_dict = {"name": name, "password": pwd}
        credentials["usernames"].update({uname: user_dict})

    
    #add a cookie which will be stored on client browser to save credentials till 30days
    authenticator = stauth.Authenticate(credentials, "pls_generator", "abcdef", cookie_expiry_days = 30)

    #u can locate the authenticator in the main body or the sidebar
    name, authentication_status, username = authenticator.login("Login", "main")
    
    if st.session_state["authentication_status"] == False:
        st.error("Username/password is incorrect")
        
    if st.session_state["authentication_status"] == None:
        st.warning("Please enter your username and password")
        
    if st.session_state["authentication_status"]:
        
        #logout button on main container
        authenticator.logout('Logout', 'main')
        st.subheader(f'Welcome *{st.session_state["name"]}*')
        
        #set bg image cover
        set_bg_hack(os.path.join(rootdir, 'iqvia-bg.jpg'))
        sidebar_bg(os.path.join(rootdir, 'iqvia-blue.png'))
        header_bg(os.path.join(rootdir, 'iqvia-dark-blue.png'))

        #setting banner image
        st.image(Image.open(os.path.join(rootdir, 'Pfizer-AI.jpg')))
        
        selected_tab = option_menu(
            menu_title=None,  # required
            options=["PLS Generator", "RCT QnA", "RCT Chatbot"],  # required
            icons=["house", "book", "envelope"],  # optional
            menu_icon="cast",  # optional
            default_index=0,  # optional
            orientation="horizontal",
        )
        
        #setting input components on sidebar
        with st.sidebar:

            st.image(Image.open(os.path.join(rootdir, 'iqvia-logo.png')))
            #setting title
            st.markdown("""<h3 style='text-align: center'>MAIA - Medical Affairs Intelligence Assistant</h3>""", unsafe_allow_html=True)

            # Step 1: Document Upload
            st.subheader("Step 1: Upload Clinical trial document")
            uploaded_file = st.file_uploader("Upload document", accept_multiple_files=False, type=["pdf"])

            # Step 2: User Inputs
            st.subheader("Step 2: Define the tone and Grade of PLS")
            # Set default values for radio button and slider
            default_tense = "Completed"
            default_pls_grade = "Low"

            # Radio button for tense selection
            tense = st.radio("Current status of the study for writing tense", options=["On-going", "Completed", "Upcoming"], key="tense", index=["On-going", "Completed", "Upcoming"].index(default_tense))
            st.write('<style>div.row-widget.stRadio > div{flex-direction:row;}</style>', unsafe_allow_html=True)

            # Slider for PLS grade selection
            #pls_grade = st.slider("Health Literacy Grade Reading level", min_value=0, max_value=10, step=5, key="pls_grade", value=default_pls_grade)
            pls_grade = st.select_slider("Health Literacy Grade of audience", options=["Low", "High"], key="pls_grade", value = default_pls_grade)

            st.session_state.process_button = False
            process_button = st.button("Process Documents")
            st.session_state.process_button = process_button
            st.session_state.uploaded_file = uploaded_file
            st.session_state.selected_tab = selected_tab
        
        #if st.session_state.process_button and st.session_state.uploaded_file:
        #if process_button and uploaded_file:
            
            # Retrieve user inputs if you haven't initialized them to any variable, then retrieve from streamlit session state
            # st.session_state.tense = tense
            # st.session_state.pls_grade = pls_grade
            
                        
        if st.session_state.uploaded_file and st.session_state.selected_tab == "RCT QnA":               
            st.subheader("Ask your PDF 💬")
            # extract the text
            if uploaded_file is not None:
              pdf_reader = PdfReader(uploaded_file)
              text = ""
              for page in pdf_reader.pages:
                text += page.extract_text()

              # split into chunks
              text_splitter = CharacterTextSplitter(
                separator="\n",
                chunk_size=1000,
                chunk_overlap=200,
                length_function=len
              )
              chunks = text_splitter.split_text(text)

              # create embeddings
              embeddings = OpenAIEmbeddings()
              knowledge_base = FAISS.from_texts(chunks, embeddings)

              # show user input
              user_question = st.text_input("Ask a question about your PDF:")
              if user_question:
                docs = knowledge_base.similarity_search(user_question)

                llm = OpenAI()
                chain = load_qa_chain(llm, chain_type="stuff")
                with get_openai_callback() as cb:
                  response = chain.run(input_documents=docs, question=user_question)
                  print(cb)

                st.write(response)    
            
        if st.session_state.selected_tab == "RCT Chatbot":    
            st.title(f"You have selected {selected_tab}")

        if st.session_state.selected_tab == "PLS Generator":
            if st.session_state.process_button and st.session_state.uploaded_file:
                col1, col2 = st.columns([0.2,0.8], gap="large")
                with col1:
                    input_file = save_uploadedfile(uploaded_file)
                    pdf_file = os.path.join(datadir, uploaded_file.name) #rootdir + "/" + uploaded_file.name
                    pdf_view = displayPDF(pdf_file)
                with col2:
                    with st.spinner(text='Processing research document you gave on the left to generate Plain Language Summary for you...⏳'):

                        # Call the processing function on the uploaded documents with user inputs
                        replacements = process_documents(pdf_file, tense, pls_grade)
                        st.success("Processed Output to be filled up in the preferred PLS template")

                        #Display processed output
                        #st.write(replacements)

                # Store the replacements dictionary in session state
                st.session_state.replacements = replacements

            # Step 3: PPT Template Selection and Download
            st.subheader("Step 3: Select PLS Template and Download")

            # Add radio buttons for template selection here    
            default_template = "Pfizer_Blue_PLS_Template"
            selected_template = st.radio("Select PPT Template", options=["Pfizer_Blue_PLS_Template", "Pfizer_Red_PLS_Template", "Pfizer_Long_PLS_Template"], index=["Pfizer_Blue_PLS_Template", "Pfizer_Red_PLS_Template", "Pfizer_Long_PLS_Template"].index(default_template))
            st.write('<style>div.row-widget.stRadio > div{flex-direction:row;}</style>', unsafe_allow_html=True)
            #selected_template = pills("", ["Pfizer_Blue_PLS_Template", "Pfizer_Red_PLS_Template", "Pfizer_Long_PLS_Template"], ["🍀", "🎈", "🌈"])

            generate_ppt_button = st.button("Generate PLS")

            if generate_ppt_button:
                # Retrieve the replacements dictionary from session state
                replacements = st.session_state.replacements
                st.session_state.process_button = False

                if replacements:
                    with st.spinner('Generating PLS slides for you...⏳'):
                        # Call the postprocessing function to generate PPT content
                        ppt_content = postprocess_to_ppt(replacements, selected_template)

                        doc_content = postprocess_to_doc(replacements)

                        # Display the PPT content using st.markdown or st.write
                        #st.markdown(ppt_content, unsafe_allow_html=True)
                        st.markdown(list(replacements.keys()))

                        # Store the PPT content in session state
                        st.session_state.ppt_content = ppt_content
                        st.session_state.doc_content = doc_content

                         # Step 4: PPT Download
                        if "ppt_content" and "doc_content" in st.session_state:
                            ppt_content = st.session_state.ppt_content
                            doc_content = st.session_state.doc_content

                            st.session_state.replacements = replacements
                            st.session_state.process_button = False

                            # Save the modified presentation object to a temporary file
                            #ppt_output_file = f"PLS_{replacements['<Title>']}_{datetime.now().strftime('%Y%m%d%H%M%S')}.pptx"                    
                            ppt_output_file = "PLS_PPT.pptx"
                            #ppt_content.save(ppt_output_file)

                            # save presentation as binary output
                            binary_output = BytesIO()
                            ppt_content.save(binary_output)

                            binary_output_doc = BytesIO()
                            doc_content.save(binary_output_doc)

                            # display success message and download button
                            st.success(':tada: The PLS template has been filled with above sections in ' + selected_template)

                            # Provide the download link for the generated PPT
                            st.download_button("Download PLS PPT", data=binary_output.getvalue(), file_name=ppt_output_file, mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
                            st.download_button("Download PLS Doc", data=binary_output_doc.getvalue(), file_name="PLS_DOC.docx", mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document")
                    
if __name__ == "__main__":
    main()
