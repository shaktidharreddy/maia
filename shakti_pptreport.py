#now building a component for the main program
#table and bar chart and combining with slide content
import streamlit as st
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_DATA_LABEL_POSITION
import os


def postprocess_to_ppt(replacements, selected_template):
    
    rootdir = "/home/cdsw/experimentation_project1/PLS_project"
    
    df = pd.read_csv(os.path.join(rootdir, 'Numericals.csv'))
    
    #selected_template = "PLS_PPT_Template"
    ppt_file = f"{selected_template}.pptx"

    # Load the existing PowerPoint presentation
    presentation = Presentation(os.path.join(rootdir, ppt_file))


    ################################################### slide content code

    # replacements = {
    #         "<Title>": """GPTAPIcall("title", tense, pls_grade)""",
    #         "<Subtitle>": """GPTAPIcall("subtitle", tense, pls_grade)""",
    #         "<Introduction>": """GPTAPIcall("introduction", tense, pls_grade)""",
    #         "<Phonetics>": 'GPTAPIcall("phonetics", tense, pls_grade)',
    #         "<Key takeaway>": 'GPTAPIcall("keytakeaway", tense, pls_grade)',
    #         "<Results>": 'GPTAPIcall("results", tense, pls_grade)',
    #         "<Intro summary>": """GPTAPIcall("conclusion", tense, pls_grade)"""
    #     }

    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        for placeholder, new_text in replacements.items():
                            if run.text == placeholder:
                                # Preserve formatting of the first run in the paragraph
                                first_run = paragraph.runs[0]
                                font_size = first_run.font.size
                                font_name = first_run.font.name
                                font_bold = first_run.font.bold
                                font_italic = first_run.font.italic

                                # Check if font color is explicitly defined
                                if first_run.font.color.type == "rgb":
                                    font_color = first_run.font.color.rgb
                                else:
                                    font_color = None

                                # Replace text while preserving formatting
                                run.text = new_text

                                # Apply preserved formatting to the entire paragraph
                                for run in paragraph.runs:
                                    run.font.size = font_size
                                    run.font.name = font_name
                                    run.font.bold = font_bold
                                    run.font.italic = font_italic
                                    if font_color:
                                        run.font.color.rgb = font_color


    ################################################### chart code

    # Find the desired text placeholder in the presentation
    placeholder_text = '<Chart>'
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                if text_frame.text == placeholder_text:
                    # Delete the existing shape (text placeholder)
                    slide.shapes._spTree.remove(shape._element)

                    # Calculate the chart position and size based on the placeholder's position and size
                    left = shape.left
                    top = shape.top
                    width = shape.width
                    height = shape.height

                    # Define the desired dimensions for the chart
                    new_width = width * 10
                    new_height = height * 20

                    # Calculate the new position for the chart to maintain the same starting position as the placeholder
                    new_left = left #- (new_width - width) / 2
                    new_top = top #- (new_height - height) / 2

                    # Create a chart on the slide with the new dimensions
                    x, y, cx, cy = new_left, new_top, new_width, new_height
                    chart_data = CategoryChartData()
                    chart_data.categories = df['Duration'].tolist()
                    chart_data.add_series('Treatment', df['Treatment'].tolist())
                    chart_data.add_series('Placebo', df['Placebo'].tolist())
                    chart = slide.shapes.add_chart(
                        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
                    ).chart

                    # Set the chart title and axis labels
                    chart.has_title = True
                    chart.chart_title.text_frame.text = 'Improvement Rate Comparison'
                    chart.chart_title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    chart.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(16)
                    chart.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black font color

                    chart.category_axis.has_major_gridlines = False
                    chart.category_axis.tick_labels.font.size = Pt(12)
                    chart.category_axis.tick_labels.font.color.rgb = RGBColor(0, 0, 0)  # Black font color

                    chart.value_axis.has_major_gridlines = True
                    chart.value_axis.major_gridlines.format.line.width = Pt(0.5)
                    chart.value_axis.major_gridlines.format.line.color.rgb = RGBColor(200, 200, 200)  # Light gray color
                    chart.value_axis.tick_labels.font.size = Pt(12)
                    chart.value_axis.tick_labels.font.color.rgb = RGBColor(0, 0, 0)  # Black font color

                    # Adjust the Y-axis scale to show up to 100%
                    chart.value_axis.maximum_scale = 100

                    # Set the legend position
                    chart.has_legend = True
                    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                    chart.legend.include_in_layout = False

                    # Apply formatting to the data series
                    series_treatment = chart.series[0]
                    fill_treatment = series_treatment.format.fill
                    fill_treatment.solid()
                    fill_treatment.fore_color.rgb = RGBColor(142, 68, 173)  # Violet shade
                    series_treatment.format.line.color.rgb = RGBColor(142, 68, 173)  # Violet shade

                    series_placebo = chart.series[1]
                    fill_placebo = series_placebo.format.fill
                    fill_placebo.solid()
                    fill_placebo.fore_color.rgb = RGBColor(79, 129, 189)  # Blue shade
                    series_placebo.format.line.color.rgb = RGBColor(79, 129, 189)  # Blue shade

                    # Apply data labels to the data series
                    series_treatment.data_labels.show_value = True
    #                 series_treatment.data_labels.number_format = '0%'
    #                 series_treatment.data_labels.position = XL_DATA_LABEL_POSITION.ABOVE
                    series_treatment.data_labels.font.size = Pt(12)
                    series_treatment.data_labels.font.color.rgb = RGBColor(0, 0, 0)  # Black font color

                    series_placebo.data_labels.show_value = True
    #                 series_placebo.data_labels.number_format = '0%'
    #                 series_placebo.data_labels.position = XL_DATA_LABEL_POSITION.ABOVE
                    series_placebo.data_labels.font.size = Pt(12)
                    series_placebo.data_labels.font.color.rgb = RGBColor(0, 0, 0)  # Black font color



    ##########################################################table code
    # Find the desired text placeholder in the presentation
    placeholder_text = '<Table>'
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                if text_frame.text == placeholder_text:
                    # Replace the text placeholder with a table using the DataFrame
                    left = shape.left
                    top = shape.top
                    width = shape.width
                    height = shape.height

                    # Delete the existing shape (text placeholder)
                    slide.shapes._spTree.remove(shape._element)

                    # Add a table to the slide using the DataFrame dimensions and position it where the placeholder was
                    rows, cols = df.shape
                    table = slide.shapes.add_table(rows=rows + 1, cols=cols, left=left, top=top, width=(width*7), height=height).table

                    # Set the column names as the header row of the table with jazzy formatting
                    header_row = table.rows[0]
                    for i, column_name in enumerate(df.columns):
                        cell = header_row.cells[i]
                        cell.text = column_name

                        # Apply formatting to the header row
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(168, 101, 168)  # Violet shade
                        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                        cell.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
                        cell.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 255, 255)  # White font color

                    # Fill in the data rows of the table with jazzy formatting
                    for i, row in enumerate(df.itertuples(), start=1):
                        for j, value in enumerate(row[1:], start=0):
                            cell = table.cell(i, j)
                            cell.text = str(value)

                            # Apply formatting to the data cells
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor(204, 153, 204)  # Light Violet shade
                            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                            cell.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
                            cell.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(64, 64, 64)  # Dark Gray font color



    return presentation

    # Save the modified presentation
    #presentation.save('modified_presentation.pptx')
